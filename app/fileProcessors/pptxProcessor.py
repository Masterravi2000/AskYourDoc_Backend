from pptx import Presentation
import os
from app.status_store import set_status 

def extract_pptx(file_path: str) -> str:
    text_parts = []
    filename = os.path.basename(file_path)

    try:
        set_status(filename, "processing")
        
        prs = Presentation(file_path)

        for slide_num, slide in enumerate(prs.slides):
            slide_text = []

            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    content = shape.text.strip()
                    if content:
                        slide_text.append(content)

            if slide_text:
                text_parts.append(
                    f"--- Slide {slide_num + 1} ---\n" + "\n".join(slide_text)
                )
                
        set_status(filename, "processed")
        
    except Exception as e:
        set_status(filename, "failed", str(e))
        print(f"{filename} → failed ❌ ({e})")
        
        if os.path.exists(file_path):
            os.remove(file_path)
            
        raise Exception(f"PPTX processing failed: {str(e)}")

    return "\n\n".join(text_parts)