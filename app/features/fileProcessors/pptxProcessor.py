from pptx import Presentation
import os
import time
from app.repositories.status_store_repository import set_status 

def extract_pptx(fileId: str, file_path: str, last_modified: float) -> str:
    documents = []
    filename = os.path.basename(file_path)
    stat  = os.stat(file_path)
    file_size = stat.st_size
    created_on = stat.st_birthtime
    
    # start = time.perf_counter()
    
    try:
        set_status(fileId, filename, "processing")
        
        prs = Presentation(file_path)

        for slide_num, slide in enumerate(prs.slides):
            slide_text = []

            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    content = shape.text.strip()
                    if content:
                        slide_text.append(content)

            if slide_text:
                documents.append({
                    "text": "\n".join(slide_text),
                    "metadata": {
                        "file_name": filename,
                        "file_type": "pptx",
                        "page": slide_num + 1,
                        "file_size": file_size,
                        "created_on": created_on,
                        "last_modified": last_modified
                    }
                })
                
        
    except Exception as e:
        set_status(fileId, filename, "failed", str(e))
        print(f"{filename} → failed ❌ ({e})")
        
        if os.path.exists(file_path):
            os.remove(file_path)
            
        raise Exception(f"PPTX processing failed: {str(e)}")
    
    # print(f"[Extraction] Completed in {time.perf_counter() - start:.3f} sec")
    
    return documents